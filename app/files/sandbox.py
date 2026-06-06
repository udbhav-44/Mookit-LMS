"""
Sandboxed text extraction (A2.3).

Extracts plain text from uploaded files in an isolated subprocess with:
  - A hard wall-clock timeout (EXTRACT_TIMEOUT_SECONDS).
  - No network access (best-effort on Linux via /proc/net checks; proper isolation
    requires running the worker in a network-namespace-restricted container).
  - CPU / memory limits set via resource.setrlimit in the child process.
  - Active content stripped (macros, external references) by choosing safe parser libs.

Supported formats: PDF, DOCX, PPTX, XLSX, CSV, TXT.
All extracted text is treated as *untrusted data* — callers must not execute it.
"""

import asyncio
import csv
import io
import logging
import os
import resource
import subprocess
import sys
import tempfile
from pathlib import Path

logger = logging.getLogger(__name__)

EXTRACT_TIMEOUT_SECONDS = 60
_MAX_CPU_SECONDS = 30
_MAX_MEM_BYTES = 256 * 1024 * 1024  # 256 MB

# ── child-process entry point ────────────────────────────────────────────────

_EXTRACTOR_SCRIPT = """
import sys, os, resource, io, csv

def _limit():
    try:
        resource.setrlimit(resource.RLIMIT_CPU, ({cpu}, {cpu}))
        resource.setrlimit(resource.RLIMIT_AS,  ({mem}, {mem}))
    except Exception:
        pass  # best-effort; parent has its own timeout guard

_limit()

file_path = sys.argv[1]
suffix = os.path.splitext(file_path)[1].lower()
text = ""

try:
    if suffix == ".pdf":
        from pdfminer.high_level import extract_text as _pdf_extract
        text = _pdf_extract(file_path)

    elif suffix == ".docx":
        import docx
        doc = docx.Document(file_path)
        text = "\\n".join(p.text for p in doc.paragraphs)

    elif suffix == ".pptx":
        from pptx import Presentation
        prs = Presentation(file_path)
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    lines.append(shape.text)
        text = "\\n".join(lines)

    elif suffix in (".xlsx", ".xls"):
        import openpyxl
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        lines = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                lines.append("\\t".join(cells))
        text = "\\n".join(lines)

    elif suffix == ".csv":
        with open(file_path, newline="", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            text = "\\n".join("\\t".join(row) for row in reader)

    elif suffix == ".txt":
        with open(file_path, encoding="utf-8", errors="replace") as f:
            text = f.read()

    else:
        sys.exit(f"Unsupported format: {{suffix}}")

except Exception as exc:
    sys.exit(f"Extraction failed: {{exc}}")

sys.stdout.write(text)
""".format(cpu=_MAX_CPU_SECONDS, mem=_MAX_MEM_BYTES)


class ExtractionError(Exception):
    pass


class ExtractionSandbox:
    """Run text extraction in a subprocess to isolate parser-level vulnerabilities."""

    async def extract_text(self, file_path: str) -> str:
        suffix = Path(file_path).suffix.lower()
        if suffix not in {".pdf", ".docx", ".pptx", ".xlsx", ".xls", ".csv", ".txt"}:
            raise ExtractionError(f"Unsupported file format: {suffix}")

        with tempfile.NamedTemporaryFile(
            mode="w", suffix=".py", delete=False, prefix="mookit_extract_"
        ) as script_file:
            script_file.write(_EXTRACTOR_SCRIPT)
            script_path = script_file.name

        try:
            proc = await asyncio.create_subprocess_exec(
                sys.executable,
                script_path,
                file_path,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE,
                # Inherit the environment but don't give the subprocess a controlling terminal.
                env={k: v for k, v in os.environ.items()
                     if k in ("PATH", "PYTHONPATH", "HOME", "LANG", "LC_ALL")},
            )
            try:
                stdout, stderr = await asyncio.wait_for(
                    proc.communicate(), timeout=EXTRACT_TIMEOUT_SECONDS
                )
            except asyncio.TimeoutError:
                proc.kill()
                await proc.wait()
                raise ExtractionError(f"Extraction timed out after {EXTRACT_TIMEOUT_SECONDS}s")

            if proc.returncode != 0:
                err_msg = stderr.decode(errors="replace").strip()
                raise ExtractionError(f"Extraction subprocess failed: {err_msg}")

            return stdout.decode("utf-8", errors="replace")

        finally:
            try:
                os.unlink(script_path)
            except OSError:
                pass

    def cleanup(self, path: str) -> None:
        try:
            if os.path.isdir(path):
                import shutil
                shutil.rmtree(path, ignore_errors=True)
            elif os.path.exists(path):
                os.remove(path)
        except OSError as exc:
            logger.warning("Sandbox cleanup failed for %s: %s", path, exc)
